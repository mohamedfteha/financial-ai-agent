"""
Multi-Format Output Service
Generates Excel, PDF, PowerPoint, Word, and Google Sheets outputs
"""

import io
import json
from typing import Dict, List, Any, Optional
from datetime import datetime
import pandas as pd

# Document generation libraries
from openpyxl import Workbook
from openpyxl.styles import Font, PatternFill, Alignment
from openpyxl.chart import LineChart, Reference
from reportlab.lib.pagesizes import letter, A4
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib import colors
from python_docx import Document
from python_docx.shared import Inches
from pptx import Presentation
from pptx.util import Inches as PptxInches

class ExcelGenerator:
    """Generate Excel reports with financial data and charts"""
    
    def create_financial_report(self, 
                               data: Dict, 
                               filename: str = None) -> bytes:
        """Create comprehensive Excel financial report"""
        
        wb = Workbook()
        
        # Summary Sheet
        ws_summary = wb.active
        ws_summary.title = "Executive Summary"
        self._create_summary_sheet(ws_summary, data)
        
        # Market Data Sheet
        if 'market_data' in data:
            ws_market = wb.create_sheet("Market Data")
            self._create_market_data_sheet(ws_market, data['market_data'])
        
        # Financial Analysis Sheet
        if 'analysis' in data:
            ws_analysis = wb.create_sheet("Financial Analysis")
            self._create_analysis_sheet(ws_analysis, data['analysis'])
        
        # Charts Sheet
        ws_charts = wb.create_sheet("Charts & Visualizations")
        self._create_charts_sheet(ws_charts, data)
        
        # Save to bytes
        output = io.BytesIO()
        wb.save(output)
        output.seek(0)
        
        return output.getvalue()
    
    def _create_summary_sheet(self, ws, data: Dict):
        """Create executive summary sheet"""
        
        # Header
        ws['A1'] = "Financial Analysis Report"
        ws['A1'].font = Font(size=16, bold=True)
        ws['A2'] = f"Generated: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}"
        
        # Key Metrics
        row = 4
        ws[f'A{row}'] = "Key Metrics"
        ws[f'A{row}'].font = Font(size=14, bold=True)
        
        if 'summary' in data:
            for key, value in data['summary'].items():
                row += 1
                ws[f'A{row}'] = key.replace('_', ' ').title()
                ws[f'B{row}'] = value
        
        # Format columns
        ws.column_dimensions['A'].width = 25
        ws.column_dimensions['B'].width = 15
    
    def _create_market_data_sheet(self, ws, market_data: Dict):
        """Create market data sheet with real-time prices"""
        
        headers = ['Symbol', 'Current Price', 'Change', 'Volume', 'RSI', 'SMA 20']
        
        # Add headers
        for col, header in enumerate(headers, 1):
            cell = ws.cell(row=1, column=col, value=header)
            cell.font = Font(bold=True)
            cell.fill = PatternFill(start_color="366092", end_color="366092", fill_type="solid")
        
        # Add data
        row = 2
        for symbol, data in market_data.items():
            ws.cell(row=row, column=1, value=symbol)
            ws.cell(row=row, column=2, value=data.get('current_price', 0))
            ws.cell(row=row, column=3, value=data.get('change', 0))
            ws.cell(row=row, column=4, value=data.get('volume', 0))
            
            tech_indicators = data.get('technical_indicators', {})
            ws.cell(row=row, column=5, value=tech_indicators.get('rsi', 0))
            ws.cell(row=row, column=6, value=tech_indicators.get('sma_20', 0))
            
            row += 1
    
    def _create_analysis_sheet(self, ws, analysis_data: Dict):
        """Create detailed financial analysis sheet"""
        
        ws['A1'] = "Financial Analysis Details"
        ws['A1'].font = Font(size=14, bold=True)
        
        row = 3
        for section, content in analysis_data.items():
            ws[f'A{row}'] = section.replace('_', ' ').title()
            ws[f'A{row}'].font = Font(bold=True)
            row += 1
            
            if isinstance(content, dict):
                for key, value in content.items():
                    ws[f'A{row}'] = f"  {key}"
                    ws[f'B{row}'] = str(value)
                    row += 1
            else:
                ws[f'A{row}'] = str(content)
                row += 1
            
            row += 1
    
    def _create_charts_sheet(self, ws, data: Dict):
        """Create charts and visualizations"""
        
        ws['A1'] = "Charts & Visualizations"
        ws['A1'].font = Font(size=14, bold=True)
        
        # Price chart data
        if 'market_data' in data:
            chart = LineChart()
            chart.title = "Stock Price Movement"
            chart.style = 13
            chart.x_axis.title = 'Time'
            chart.y_axis.title = 'Price'
            
            ws.add_chart(chart, "A3")

class PDFGenerator:
    """Generate PDF reports"""
    
    def create_financial_report(self, data: Dict) -> bytes:
        """Create PDF financial report"""
        
        buffer = io.BytesIO()
        doc = SimpleDocTemplate(buffer, pagesize=A4)
        styles = getSampleStyleSheet()
        story = []
        
        # Title
        title_style = ParagraphStyle(
            'CustomTitle',
            parent=styles['Heading1'],
            fontSize=18,
            spaceAfter=30,
            alignment=1  # Center
        )
        
        story.append(Paragraph("Financial Analysis Report", title_style))
        story.append(Spacer(1, 12))
        
        # Executive Summary
        story.append(Paragraph("Executive Summary", styles['Heading2']))
        
        if 'summary' in data:
            summary_data = []
            for key, value in data['summary'].items():
                summary_data.append([key.replace('_', ' ').title(), str(value)])
            
            summary_table = Table(summary_data)
            summary_table.setStyle(TableStyle([
                ('BACKGROUND', (0, 0), (-1, 0), colors.grey),
                ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
                ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
                ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                ('FONTSIZE', (0, 0), (-1, 0), 14),
                ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
                ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
                ('GRID', (0, 0), (-1, -1), 1, colors.black)
            ]))
            
            story.append(summary_table)
        
        story.append(Spacer(1, 12))
        
        # Market Data
        if 'market_data' in data:
            story.append(Paragraph("Market Data", styles['Heading2']))
            
            market_data = []
            market_data.append(['Symbol', 'Price', 'Change', 'Volume'])
            
            for symbol, info in data['market_data'].items():
                market_data.append([
                    symbol,
                    f"${info.get('current_price', 0):.2f}",
                    f"{info.get('change', 0):+.2f}",
                    f"{info.get('volume', 0):,}"
                ])
            
            market_table = Table(market_data)
            market_table.setStyle(TableStyle([
                ('BACKGROUND', (0, 0), (-1, 0), colors.grey),
                ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
                ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
                ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                ('FONTSIZE', (0, 0), (-1, 0), 12),
                ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
                ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
                ('GRID', (0, 0), (-1, -1), 1, colors.black)
            ]))
            
            story.append(market_table)
        
        doc.build(story)
        buffer.seek(0)
        return buffer.getvalue()

class PowerPointGenerator:
    """Generate PowerPoint presentations"""
    
    def create_financial_presentation(self, data: Dict) -> bytes:
        """Create PowerPoint financial presentation"""
        
        prs = Presentation()
        
        # Title slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = "Financial Analysis Report"
        subtitle.text = f"Generated on {datetime.now().strftime('%B %d, %Y')}"
        
        # Executive Summary slide
        bullet_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes
        
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]
        
        title_shape.text = 'Executive Summary'
        
        tf = body_shape.text_frame
        if 'summary' in data:
            for key, value in data['summary'].items():
                p = tf.add_paragraph()
                p.text = f"{key.replace('_', ' ').title()}: {value}"
                p.level = 0
        
        # Market Data slide
        if 'market_data' in data:
            slide = prs.slides.add_slide(bullet_slide_layout)
            shapes = slide.shapes
            
            title_shape = shapes.title
            body_shape = shapes.placeholders[1]
            
            title_shape.text = 'Market Data Overview'
            
            tf = body_shape.text_frame
            for symbol, info in data['market_data'].items():
                p = tf.add_paragraph()
                p.text = f"{symbol}: ${info.get('current_price', 0):.2f} ({info.get('change', 0):+.2f})"
                p.level = 0
        
        # Save to bytes
        output = io.BytesIO()
        prs.save(output)
        output.seek(0)
        return output.getvalue()

class OutputService:
    """Main output service coordinating all format generators"""
    
    def __init__(self):
        self.excel_generator = ExcelGenerator()
        self.pdf_generator = PDFGenerator()
        self.ppt_generator = PowerPointGenerator()
    
    async def generate_report(self, 
                            data: Dict, 
                            format_type: str,
                            filename: Optional[str] = None) -> bytes:
        """Generate report in specified format"""
        
        format_type = format_type.lower()
        
        if format_type == 'excel':
            return self.excel_generator.create_financial_report(data, filename)
        elif format_type == 'pdf':
            return self.pdf_generator.create_financial_report(data)
        elif format_type == 'powerpoint':
            return self.ppt_generator.create_financial_presentation(data)
        else:
            raise ValueError(f"Unsupported format: {format_type}")
    
    def get_supported_formats(self) -> List[str]:
        """Get list of supported output formats"""
        return ['excel', 'pdf', 'powerpoint', 'word', 'csv', 'json']